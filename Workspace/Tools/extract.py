from pptx import Presentation
import sys
p = Presentation(sys.argv[1])
for i, slide in enumerate(p.slides, 1):
    print(f"\n=== Slide {i} ===")
    for shape in slide.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if t: print(t)
        elif shape.shape_type == 13:
            print(f"[picture: {getattr(shape.image,'filename',None) or shape.name}]")
    if slide.has_notes_slide:
        n = slide.notes_slide.notes_text_frame.text.strip()
        if n: print(f"--- notes: {n}")
