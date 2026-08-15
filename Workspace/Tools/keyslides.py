from pptx import Presentation
import sys, re
KEY = re.compile(r'lesson|takeaway|summary|key point|warning|advice|main strateg|important|remember', re.I)
for f in sys.argv[1:]:
    p = Presentation(f)
    print(f"\n######## {f.split('/')[-1]}")
    for i, s in enumerate(p.slides,1):
        txts=[]
        for sh in s.shapes:
            if sh.has_text_frame:
                t=sh.text_frame.text.strip()
                if t: txts.append(t)
        full="\n".join(txts)
        if KEY.search(full):
            print(f"--- slide {i}:")
            print(full[:500])
