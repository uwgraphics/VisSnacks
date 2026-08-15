from pptx import Presentation
import sys, glob
for f in sys.argv[1:]:
    try:
        p = Presentation(f)
        print(f"\n######## {f.split('/')[-1]} ({len(p.slides.__iter__.__self__._sldIdLst)} slides)")
        for i, s in enumerate(p.slides,1):
            t=""
            try:
                if s.shapes.title and s.shapes.title.has_text_frame:
                    t = s.shapes.title.text.strip().replace("\n"," / ")
            except Exception: pass
            if t: print(f"{i}: {t}")
    except Exception as e:
        print(f"\n######## {f}: ERROR {e}")
